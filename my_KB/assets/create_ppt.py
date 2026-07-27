"""
Run this script to generate rolling-summarization-plan.pptx
Requires: pip install python-pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

ACCENT  = RGBColor(0x00, 0x70, 0xC0)   # blue
GREEN   = RGBColor(0x00, 0x70, 0x50)
RED     = RGBColor(0xC0, 0x00, 0x00)
DARK    = RGBColor(0x1F, 0x1F, 0x1F)
LIGHT   = RGBColor(0xF5, 0xF5, 0xF5)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0x60, 0x60, 0x60)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]  # completely blank


def bg(slide, color=LIGHT):
    from pptx.util import Pt
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, l, t, w, h, text, font_size=18, bold=False,
        color=DARK, bg_color=None, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def rect(slide, l, t, w, h, fill_color, line_color=None):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


# ── Slide 1: Title ──────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(blank)
bg(s1, RGBColor(0x1A, 0x1A, 0x2E))
rect(s1, 0, 2.8, 13.33, 0.06, ACCENT)
box(s1, 1, 1.2, 11, 1.2,
    "Multi-Turn Conversation Memory", 40, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER)
box(s1, 1, 2.5, 11, 0.6,
    "Rolling LLM Summarization — AgentMesh", 20,
    color=RGBColor(0xAA, 0xCC, 0xFF), align=PP_ALIGN.CENTER)
box(s1, 1, 4.0, 11, 0.5,
    "Replace hard turn-limit (3 turns) with an always-current running summary", 16,
    color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)
box(s1, 1, 5.0, 11, 0.4,
    "Stack: Microsoft Agent Framework · Groq (OpenAI-compat) · JSONL session store", 13,
    color=GRAY, align=PP_ALIGN.CENTER)


# ── Slide 2: Problem ────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(blank)
bg(s2)
rect(s2, 0, 0, 13.33, 0.9, ACCENT)
box(s2, 0.3, 0.1, 12, 0.7, "The Problem", 28, bold=True, color=WHITE)

problems = [
    ("Hard cap at 3 turns",       "CONVERSATION_MAX_TURNS=3  →  last 6 messages sent to LLM", RED),
    ("Silent context loss",       "Turns older than 3 are dropped. LLM has no memory of them.", RED),
    ("Poor long-session quality", "User must re-explain context that already happened.", RED),
]
for i, (title, desc, col) in enumerate(problems):
    y = 1.2 + i * 1.5
    rect(s2, 0.5, y, 12.3, 1.2, WHITE, ACCENT)
    box(s2, 0.8, y + 0.05, 12, 0.45, f"✗  {title}", 16, bold=True, color=col)
    box(s2, 0.8, y + 0.5,  12, 0.55, desc, 13, color=DARK)

box(s2, 0.5, 5.7, 12.3, 0.4,
    "Code location:  src/config.py:48   ·   src/mesh/orchestrator.py:114   ·   src/memory/conversation_store.py:47",
    11, color=GRAY)


# ── Slide 3: Goal ───────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(blank)
bg(s3)
rect(s3, 0, 0, 13.33, 0.9, RGBColor(0x00, 0x70, 0x50))
box(s3, 0.3, 0.1, 12, 0.7, "The Goal", 28, bold=True, color=WHITE)

goals = [
    "No hard turn limit — unlimited conversation length",
    "Summarize all prior turns into a ≤200-word rolling summary",
    "Always send  [Summary] + [Current question]  to the LLM — every turn",
    "Persist summary in JSONL session file — survives server restarts",
    "Non-blocking — summarization fires async after response is returned",
]
for i, g in enumerate(goals):
    box(s3, 0.9, 1.2 + i * 1.0, 11.5, 0.8, f"✓  {g}", 15, color=DARK)


# ── Slide 4: Options ────────────────────────────────────────────────────────
s4 = prs.slides.add_slide(blank)
bg(s4)
rect(s4, 0, 0, 13.33, 0.9, DARK)
box(s4, 0.3, 0.1, 12, 0.7, "Options Evaluated", 28, bold=True, color=WHITE)

cols = [
    ("Option A\nSemantic Kernel\nConversationSummaryMemory",
     "Native SK plugin.\nNot viable — project uses Agent Framework, not SK.",
     RED, "✗ Not viable"),
    ("Option B\nRolling LLM\nSummarization",
     "LLM call after each turn.\nStores summary in JSONL.\nSame Groq endpoint.",
     GREEN, "✓ Recommended"),
    ("Option C\nHeuristic /\nExtractive",
     "Keyword extraction, no LLM.\nSimpler but low quality.\nMisses intent.",
     RGBColor(0xB0, 0x70, 0x00), "✗ Not recommended"),
]
for i, (title, desc, col, verdict) in enumerate(cols):
    x = 0.5 + i * 4.25
    rect(s4, x, 1.1, 4.0, 5.5, WHITE, col)
    box(s4, x + 0.15, 1.2,  3.7, 1.2, title,   14, bold=True, color=col)
    box(s4, x + 0.15, 2.5,  3.7, 3.0, desc,    13, color=DARK)
    box(s4, x + 0.15, 5.7,  3.7, 0.6, verdict, 13, bold=True, color=col)


# ── Slide 5: Flow Diagram ───────────────────────────────────────────────────
s5 = prs.slides.add_slide(blank)
bg(s5)
rect(s5, 0, 0, 13.33, 0.9, ACCENT)
box(s5, 0.3, 0.1, 12, 0.7, "Per-Turn Flow (Option B)", 28, bold=True, color=WHITE)

steps = [
    ("1", "load_with_summary\n(session_id)",       "Returns (summary_str, messages[])"),
    ("2", "Build prompt\n[Summary]+[Question]",    "Compact context block sent to LLM"),
    ("3", "Call LLM via A2A\n(PriceAssistAgent)", "Answer returned to user immediately"),
    ("4", "async summarize\n& persist",            "Non-blocking — no latency impact"),
]
for i, (num, title, sub) in enumerate(steps):
    x = 0.4 + i * 3.15
    rect(s5, x, 1.2, 2.7, 1.4, ACCENT)
    box(s5, x + 0.1, 1.25, 2.5, 0.5, num, 22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(s5, x + 0.1, 1.7,  2.5, 0.8, title, 12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(s5, x + 0.1, 2.8,  2.5, 0.7, sub, 11, color=DARK)
    if i < 3:
        box(s5, x + 2.75, 1.6, 0.35, 0.6, "→", 22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

box(s5, 0.5, 3.8, 12.3, 0.4,
    "Prompt format every turn:", 13, bold=True, color=DARK)
rect(s5, 0.5, 4.3, 12.3, 2.4, RGBColor(0xE8, 0xF0, 0xFE))
box(s5, 0.7, 4.4, 11.9, 2.2,
    "[Conversation Summary]\n<running summary of all prior turns — ≤200 words>\n\n[Current question]\n<user message>",
    13, color=DARK)


# ── Slide 6: Files to Change ─────────────────────────────────────────────────
s6 = prs.slides.add_slide(blank)
bg(s6)
rect(s6, 0, 0, 13.33, 0.9, DARK)
box(s6, 0.3, 0.1, 12, 0.7, "Files to Modify", 28, bold=True, color=WHITE)

rows = [
    ("src/memory/summarizer.py",        "NEW",     "LLM summarization client — same Groq endpoint, async"),
    ("src/memory/base.py",              "Extend",  "Add abstract load_summary / save_summary methods"),
    ("src/memory/jsonl_backend.py",     "Extend",  "Implement summary read/write on JSONL file"),
    ("src/memory/conversation_store.py","Extend",  "Add load_with_summary() and save_summary() facade"),
    ("src/mesh/workflow.py",            "Modify",  "Swap format_history_block → format_summary_block; fire async summarization"),
    ("src/mesh/orchestrator.py",        "Modify",  "Line 114: use load_with_summary instead of load"),
    ("src/config.py",                   "Modify",  "Deprecate CONVERSATION_MAX_TURNS; add optional SUMMARY_MODEL"),
]
colors = {"NEW": GREEN, "Extend": ACCENT, "Modify": RGBColor(0xB0, 0x70, 0x00)}
for i, (file, change, detail) in enumerate(rows):
    y = 1.1 + i * 0.75
    rect(s6, 0.4, y, 12.5, 0.65, WHITE if i % 2 == 0 else RGBColor(0xF0, 0xF0, 0xF0))
    col = colors.get(change, DARK)
    box(s6, 0.5,  y + 0.05, 3.8, 0.55, file,   12, bold=True, color=DARK)
    box(s6, 4.3,  y + 0.05, 1.1, 0.55, change, 12, bold=True, color=col)
    box(s6, 5.5,  y + 0.05, 7.2, 0.55, detail, 12, color=GRAY)


# ── Slide 7: Verification ───────────────────────────────────────────────────
s7 = prs.slides.add_slide(blank)
bg(s7)
rect(s7, 0, 0, 13.33, 0.9, RGBColor(0x00, 0x70, 0x50))
box(s7, 0.3, 0.1, 12, 0.7, "Verification Checklist", 28, bold=True, color=WHITE)

checks = [
    "Ask 5+ questions in one session — confirm no 3-turn cap, all context preserved via summary",
    "Inspect data/conversations/{session_id}.jsonl — confirm type=summary records appear after each turn",
    "Kill + restart server mid-session — confirm summary persists and is used on next turn",
    "Confirm strip_history_echo() still works  (new [Conversation Summary] header is distinct from old [Conversation so far])",
    "Measure response latency — summarization must NOT block the response path (async task)",
]
for i, c in enumerate(checks):
    y = 1.2 + i * 1.0
    rect(s7, 0.5, y, 0.55, 0.55, ACCENT)
    box(s7, 0.55, y + 0.02, 0.45, 0.5, "☐", 16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(s7, 1.2,  y, 11.5, 0.7, c, 13, color=DARK)


# ── Save ─────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "rolling-summarization-plan.pptx")
prs.save(out)
print(f"Saved: {out}")
